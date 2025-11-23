"""
app/services/report_generation/ppt_generator.py
Generación de presentaciones PowerPoint personalizadas
"""

import os
from typing import Dict, List, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.models.student_profile import StudentProfile
from app.services.report_generation.data_visualizer import data_visualizer
from app.services.ai.gemini_service import gemini_service
from app.utils.logger import logger


class PPTGenerator:
    """
    Generador de presentaciones PowerPoint personalizadas
    Adapta contenido según el perfil del estudiante
    """
    
    def __init__(self):
        self.logger = logger
        # Directorio de salida
        self.output_dir = os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
            'generated', 'ppt'
        )
        os.makedirs(self.output_dir, exist_ok=True)
    
    # =====================================================
    # MÉTODO PRINCIPAL
    # =====================================================
    
    def generate_student_report_ppt(
        self, 
        profile: StudentProfile,
        user_id: int,
        title: str = "Reporte de Rendimiento Académico",
        include_recommendations: bool = True
    ) -> Dict:
        """
        Genera presentación completa del reporte del estudiante
        
        Args:
            profile: Perfil del estudiante
            user_id: ID del usuario
            title: Título de la presentación
            include_recommendations: Incluir slide de recomendaciones
            
        Returns:
            Dict con ruta del archivo y metadata
        """
        try:
            self.logger.info(f"Generando PPT para user_id={user_id}")
            
            # Crear presentación
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Portada
            self._add_title_slide(prs, title, profile)
            
            # Slide 2: Resumen Ejecutivo
            self._add_executive_summary_slide(prs, profile)
            
            # Slide 3: Métricas Clave
            self._add_key_metrics_slide(prs, profile)
            
            # Slide 4: Preparación para Tesis
            self._add_thesis_readiness_slide(prs, profile)
            
            # Slide 5: Fortalezas
            self._add_strengths_slide(prs, profile)
            
            # Slide 6: Áreas de Mejora
            self._add_weaknesses_slide(prs, profile)
            
            # Slide 7: Recomendaciones (opcional)
            if include_recommendations:
                self._add_recommendations_slide(prs, profile)
            
            # Slide 8: Próximos Pasos
            self._add_next_steps_slide(prs, profile, user_id)
            
            # Guardar archivo
            filename = f"reporte_{user_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)
            
            self.logger.info(f"PPT generado: {filepath}")
            
            return {
                'success': True,
                'filepath': filepath,
                'filename': filename,
                'slides_count': len(prs.slides),
                'file_size': os.path.getsize(filepath)
            }
            
        except Exception as e:
            self.logger.error(f"Error generando PPT: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }
    
    # =====================================================
    # SLIDES INDIVIDUALES
    # =====================================================
    
    def _add_title_slide(self, prs: Presentation, title: str, profile: StudentProfile):
        """Slide de portada"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Fondo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(41, 128, 185)  # Azul
        
        # Título principal
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtítulo con nombre del estudiante
        from app.models.user import User
        user = User.query.get(profile.user_id)
        student_name = f"{user.first_name} {user.last_name}" if user else "Estudiante"
        
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(8), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = student_name
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Fecha
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%d de %B, %Y")
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = RGBColor(255, 255, 255)
        date_para.alignment = PP_ALIGN.CENTER
    
    def _add_executive_summary_slide(self, prs: Presentation, profile: StudentProfile):
        """Slide de resumen ejecutivo"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        # Título
        title = slide.shapes.title
        title.text = "Resumen Ejecutivo"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(41, 128, 185)
        
        # Contenido - Resumen IA
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8), Inches(8), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Usar resumen IA del perfil
        summary_text = profile.ai_profile_summary or "Perfil en construcción."
        
        # Dividir en párrafos
        paragraphs = summary_text.split('\n\n')
        
        for i, para_text in enumerate(paragraphs[:3]):  # Máximo 3 párrafos
            if i > 0:
                text_frame.add_paragraph()
            
            para = text_frame.paragraphs[i]
            para.text = para_text.strip()
            para.font.size = Pt(14)
            para.font.name = 'Calibri'
            para.space_after = Pt(12)
            para.line_spacing = 1.2
    
    def _add_key_metrics_slide(self, prs: Presentation, profile: StudentProfile):
        """Slide de métricas clave con números grandes"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Título
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Métricas Clave"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(41, 128, 185)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Obtener stats
        stats = data_visualizer.get_summary_stats(profile, profile.user_id)
        
        # Layout: 2 filas x 3 columnas
        metrics = [
            {
                'value': stats['total_documents'],
                'label': 'Documentos\nAnalizados',
                'color': RGBColor(52, 152, 219)
            },
            {
                'value': stats['total_sessions'],
                'label': 'Sesiones\nCompletadas',
                'color': RGBColor(46, 204, 113)
            },
            {
                'value': f"{stats['thesis_score']:.0f}",
                'label': 'Score Preparación\nTesis',
                'color': RGBColor(155, 89, 182)
            },
            {
                'value': f"{stats['writing_quality']:.0f}",
                'label': 'Calidad de\nEscritura',
                'color': RGBColor(230, 126, 34)
            },
            {
                'value': f"{stats['vocabulary_score']:.0f}",
                'label': 'Riqueza de\nVocabulario',
                'color': RGBColor(231, 76, 60)
            },
            {
                'value': stats['attention_span'],
                'label': 'Span Atención\n(minutos)',
                'color': RGBColor(26, 188, 156)
            }
        ]
        
        # Posiciones
        positions = [
            (1.5, 2, 2, 1.5),    # Top left
            (4, 2, 2, 1.5),      # Top center
            (6.5, 2, 2, 1.5),    # Top right
            (1.5, 4, 2, 1.5),    # Bottom left
            (4, 4, 2, 1.5),      # Bottom center
            (6.5, 4, 2, 1.5)     # Bottom right
        ]
        
        for i, metric in enumerate(metrics):
            left, top, width, height = positions[i]
            
            # Box del número
            value_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height * 0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = str(metric['value'])
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(48)
            value_para.font.bold = True
            value_para.font.color.rgb = metric['color']
            value_para.alignment = PP_ALIGN.CENTER
            
            # Box del label
            label_box = slide.shapes.add_textbox(
                Inches(left), Inches(top + height * 0.6), Inches(width), Inches(height * 0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = metric['label']
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(12)
            label_para.font.color.rgb = RGBColor(100, 100, 100)
            label_para.alignment = PP_ALIGN.CENTER
    
    def _add_thesis_readiness_slide(self, prs: Presentation, profile: StudentProfile):
        """Slide de preparación para tesis"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = "Preparación para Tesis"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(41, 128, 185)
        
        # Score grande
        score_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(3.5), Inches(2)
        )
        score_frame = score_box.text_frame
        score_frame.text = f"{float(profile.thesis_readiness_score or 0):.0f}/100"
        score_para = score_frame.paragraphs[0]
        score_para.font.size = Pt(72)
        score_para.font.bold = True
        
        # Color según nivel
        level = profile.thesis_readiness_level
        if level == 'excelente':
            color = RGBColor(46, 204, 113)  # Verde
        elif level == 'alto':
            color = RGBColor(52, 152, 219)  # Azul
        elif level == 'medio':
            color = RGBColor(241, 196, 15)  # Amarillo
        else:
            color = RGBColor(231, 76, 60)   # Rojo
        
        score_para.font.color.rgb = color
        score_para.alignment = PP_ALIGN.CENTER
        
        # Nivel
        level_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(3.5), Inches(0.5)
        )
        level_frame = level_box.text_frame
        level_text = {
            'bajo': 'Nivel Bajo',
            'medio': 'Nivel Medio',
            'alto': 'Nivel Alto',
            'excelente': 'Nivel Excelente'
        }
        level_frame.text = level_text.get(level, 'No Evaluado')
        level_para = level_frame.paragraphs[0]
        level_para.font.size = Pt(24)
        level_para.font.bold = True
        level_para.font.color.rgb = color
        level_para.alignment = PP_ALIGN.CENTER
        
        # Detalles a la derecha
        details_box = slide.shapes.add_textbox(
            Inches(5), Inches(2), Inches(4), Inches(3.5)
        )
        details_frame = details_box.text_frame
        details_frame.word_wrap = True
        
        # Factores
        factors_text = f"""Factores Evaluados:

• Documentos analizados: {profile.total_documents_analyzed}
• Calidad de escritura: {float(profile.avg_writing_quality or 0):.1f}/100
• Vocabulario técnico: {float(profile.avg_vocabulary_richness or 0):.1f}/100
• Consistencia: {profile.writing_improvement_trend or 'Sin datos'}

Tiempo estimado de preparación:
{profile.estimated_preparation_months or 12} meses"""
        
        details_frame.text = factors_text
        for para in details_frame.paragraphs:
            para.font.size = Pt(14)
            para.space_after = Pt(6)
    
    def _add_strengths_slide(self, prs: Presentation, profile: StudentProfile):
        """Slide de fortalezas"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = "💪 Fortalezas Identificadas"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)
        
        # Recopilar todas las fortalezas
        all_strengths = []
        if profile.academic_strengths:
            all_strengths.extend(profile.academic_strengths)
        if profile.writing_strengths:
            all_strengths.extend(profile.writing_strengths)
        if profile.technical_strengths:
            all_strengths.extend(profile.technical_strengths)
        
        if not all_strengths:
            all_strengths = ["Perfil en construcción - aún identificando fortalezas"]
        
        # Contenido
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2), Inches(7), Inches(4.5)
        )
        text_frame = content_box.text_frame
        
        for strength in all_strengths[:6]:  # Máximo 6
            para = text_frame.add_paragraph()
            para.text = f"✓ {strength}"
            para.font.size = Pt(18)
            para.space_before = Pt(12)
            para.level = 0
    
    def _add_weaknesses_slide(self, prs: Presentation, profile: StudentProfile):
        """Slide de áreas de mejora"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = "🎯 Áreas de Oportunidad"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 126, 34)
        
        # Recopilar debilidades
        all_weaknesses = []
        if profile.academic_weaknesses:
            all_weaknesses.extend(profile.academic_weaknesses)
        if profile.writing_weaknesses:
            all_weaknesses.extend(profile.writing_weaknesses)
        if profile.areas_for_improvement:
            all_weaknesses.extend(profile.areas_for_improvement)
        
        if not all_weaknesses:
            all_weaknesses = ["Sin áreas críticas identificadas"]
        
        # Contenido
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2), Inches(7), Inches(4.5)
        )
        text_frame = content_box.text_frame
        
        for weakness in all_weaknesses[:6]:
            para = text_frame.add_paragraph()
            para.text = f"→ {weakness}"
            para.font.size = Pt(18)
            para.space_before = Pt(12)
            para.level = 0
    
    def _add_recommendations_slide(self, prs: Presentation, profile: StudentProfile):
        """Slide de recomendaciones"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = "💡 Recomendaciones Personalizadas"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)
        
        # Recomendaciones
        all_recs = []
        if profile.study_recommendations:
            all_recs.extend(profile.study_recommendations)
        if profile.resource_recommendations:
            all_recs.extend(profile.resource_recommendations)
        
        if not all_recs:
            all_recs = ["Continuar usando la plataforma para obtener recomendaciones personalizadas"]
        
        # Contenido
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2), Inches(7), Inches(4.5)
        )
        text_frame = content_box.text_frame
        
        for i, rec in enumerate(all_recs[:5], 1):
            para = text_frame.add_paragraph()
            para.text = f"{i}. {rec}"
            para.font.size = Pt(16)
            para.space_before = Pt(10)
            para.level = 0
    
    def _add_next_steps_slide(self, prs: Presentation, profile: StudentProfile, user_id: int):
        """Slide de próximos pasos"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = "🚀 Próximos Pasos"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(41, 128, 185)
        
        # Generar pasos con IA
        next_steps = self._generate_next_steps_with_ai(profile, user_id)
        
        # Contenido
        content_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2), Inches(7), Inches(4.5)
        )
        text_frame = content_box.text_frame
        
        for i, step in enumerate(next_steps[:4], 1):
            para = text_frame.add_paragraph()
            para.text = f"{i}. {step}"
            para.font.size = Pt(18)
            para.space_before = Pt(15)
            para.level = 0
    
    def _generate_next_steps_with_ai(self, profile: StudentProfile, user_id: int) -> List[str]:
        """Genera próximos pasos con IA"""
        try:
            prompt = f"""Basado en este perfil académico, genera 4 próximos pasos ESPECÍFICOS y ACCIONABLES para el estudiante:

Nivel de preparación tesis: {profile.thesis_readiness_level}
Score: {float(profile.thesis_readiness_score or 0):.0f}/100
Documentos: {profile.total_documents_analyzed}
Sesiones: {profile.total_sessions_completed}

Genera SOLO 4 pasos, uno por línea, sin numeración, sin formato markdown."""
            
            response = gemini_service.generate_content(prompt, user_id, "next_steps")
            
            if response['success']:
                steps = [s.strip() for s in response['content'].split('\n') if s.strip()]
                return steps[:4]
            
        except:
            pass
        
        # Fallback
        return [
            "Continuar subiendo documentos académicos para análisis",
            "Completar al menos 2 sesiones de estudio semanales",
            "Revisar y aplicar las recomendaciones personalizadas",
            "Actualizar el perfil mensualmente para seguir el progreso"
        ]


# Instancia única
ppt_generator = PPTGenerator()